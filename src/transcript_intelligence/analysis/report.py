"""
Builds the auto-generated .pptx slide deck using python-pptx.

Slide structure:
  1.  Title
  2.  Dataset Overview & Methodology
  3.  Topic Landscape (cluster map)
  4.  Topic Clusters — Key Findings
  5.  Sentiment by Call Type
  6.  Sentiment Trends Over Time
  7.  Churn Risk Analysis
  8.  Escalation Detection
  9.  Action Item Ownership
  10. Speaker Engagement
  11. Architecture & Approach
  12. Recommendations
"""
from __future__ import annotations

from pathlib import Path
from statistics import mean

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from transcript_intelligence.models import (
    ChurnSignal,
    ClusterResult,
    EscalationFlag,
    PipelineState,
    SentimentAggregate,
    TranscriptRecord,
)

# Aegis-style palette
DARK = RGBColor(0x1A, 0x1A, 0x2E)     # dark navy
ACCENT = RGBColor(0x4C, 0x9B, 0xE8)   # blue
ACCENT2 = RGBColor(0x6B, 0xCB, 0x77)  # green
WARN = RGBColor(0xE8, 0x6B, 0x4C)     # red-orange
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6C, 0x75, 0x7D)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def build_pptx(state: PipelineState, chart_paths: list[str], out_dir: Path) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    transcripts = state.get("classified_transcripts", [])
    cluster_results = state.get("cluster_results", [])
    churn_signals = state.get("churn_signals", [])
    escalation_flags = state.get("escalation_flags", [])
    action_items = state.get("action_item_ownership", [])
    speaker_engagements = state.get("speaker_engagements", [])
    sentiment_aggregates = state.get("sentiment_aggregates", [])
    stats = state.get("classification_stats", {})

    chart_map = {Path(p).stem: p for p in chart_paths}

    _slide_title(prs, len(transcripts))
    _slide_overview(prs, transcripts, cluster_results, stats, state)
    _slide_cluster_map(prs, chart_map)
    _slide_cluster_details(prs, cluster_results)
    _slide_sentiment_box(prs, transcripts, sentiment_aggregates, chart_map)
    _slide_sentiment_trend(prs, chart_map)
    _slide_churn_risk(prs, churn_signals, chart_map)
    _slide_escalation(prs, escalation_flags, chart_map)
    _slide_action_items(prs, action_items, chart_map)
    _slide_speaker_engagement(prs, speaker_engagements, transcripts, chart_map)
    _slide_architecture(prs)
    _slide_recommendations(prs, churn_signals, escalation_flags, transcripts)

    out_path = out_dir / "transcript_intelligence.pptx"
    prs.save(str(out_path))
    return str(out_path)


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation) -> object:
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def _add_bg(slide, color: RGBColor = LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, title: str, subtitle: str = ""):
    """Dark header bar across the top of the slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, Inches(1.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.25)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.alignment = PP_ALIGN.LEFT


def _add_text_box(slide, text: str, left, top, width, height,
                  size=Pt(11), bold=False, color=DARK, align=PP_ALIGN.LEFT,
                  wrap=True) -> object:
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txb


def _add_image_safe(slide, path_or_key, chart_map: dict, left, top, width, height):
    actual = chart_map.get(path_or_key) if path_or_key in chart_map else path_or_key
    if actual and Path(actual).exists():
        slide.shapes.add_picture(actual, left, top, width, height)
    else:
        # Placeholder box
        ph = slide.shapes.add_shape(1, left, top, width, height)
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        ph.text_frame.paragraphs[0].text = f"[Chart: {path_or_key}]"


def _stat_box(slide, label: str, value: str, left, top, color: RGBColor = ACCENT):
    box = slide.shapes.add_shape(1, left, top, Inches(2.5), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    tf = box.text_frame
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.text = value
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
    p2.alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# Individual slides
# ---------------------------------------------------------------------------

def _slide_title(prs, n_transcripts: int):
    slide = _blank_slide(prs)
    _add_bg(slide, DARK)

    # Title
    _add_text_box(
        slide, "Transcript Intelligence",
        Inches(1), Inches(1.8), Inches(11), Inches(1.5),
        size=Pt(48), bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide, "Insights from 100 B2B SaaS Call Transcripts",
        Inches(1), Inches(3.2), Inches(11), Inches(0.8),
        size=Pt(20), color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide, "AI Engineer Assignment  •  May 2026",
        Inches(1), Inches(5.5), Inches(11), Inches(0.6),
        size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER,
    )
    # Accent line
    line = slide.shapes.add_shape(1, Inches(4.5), Inches(4.0), Inches(4.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def _slide_overview(prs, transcripts, cluster_results, stats, state):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Dataset Overview", "What we analyzed and how")

    n = len(transcripts)
    n_internal = stats.get("internal", 0)
    n_support = stats.get("support", 0)
    n_external = stats.get("external", 0)
    n_clusters = state.get("num_clusters", 0)

    _stat_box(slide, "Total Transcripts", str(n), Inches(0.4), Inches(1.5), ACCENT)
    _stat_box(slide, "Internal Calls", str(n_internal), Inches(3.1), Inches(1.5), DARK)
    _stat_box(slide, "Support Calls", str(n_support), Inches(5.8), Inches(1.5), WARN)
    _stat_box(slide, "External Calls", str(n_external), Inches(8.5), Inches(1.5), ACCENT2)
    _stat_box(slide, "Topic Clusters", str(n_clusters), Inches(11.2), Inches(1.5), RGBColor(0x8E, 0x44, 0xAD))

    methodology = (
        "Approach:\n"
        "• Call type classification: hybrid rule-based (email domain + title keywords) + "
        "Claude claude-haiku-4-5 fallback\n"
        "• Topic clustering: sentence-transformers embeddings (all-MiniLM-L6-v2) → "
        "UMAP (10D) → HDBSCAN → Claude claude-sonnet-4-6 cluster labeling\n"
        "• Sentiment: aggregated from pre-extracted sentence-level sentimentType + summary sentimentScore\n"
        "• Bonus insights: Churn Risk Scoring, Action Item Ownership, "
        "Escalation Detection, Speaker Engagement Analysis\n"
        "• Pipeline orchestrated with LangGraph; all code in Python 3.12 + uv"
    )
    _add_text_box(
        slide, methodology,
        Inches(0.4), Inches(3.0), Inches(12.5), Inches(4.0),
        size=Pt(11), color=DARK,
    )


def _slide_cluster_map(prs, chart_map: dict):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Topic Landscape", "All 100 transcripts plotted by semantic similarity (UMAP 2D)")

    _add_image_safe(slide, "cluster_map", chart_map,
                    Inches(0.3), Inches(1.3), Inches(8.5), Inches(5.8))

    _add_text_box(
        slide,
        "Each point = one transcript\n"
        "Color = topic cluster\n"
        "Shape = call type\n"
        "Size = sentiment score",
        Inches(9.1), Inches(2.0), Inches(4.0), Inches(3.0),
        size=Pt(11), color=DARK,
    )


def _slide_cluster_details(prs, cluster_results: list[ClusterResult]):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Topic Clusters — Key Findings",
                    "Identified by HDBSCAN + labeled by Claude claude-sonnet-4-6")

    y = Inches(1.4)
    col1_x = Inches(0.4)
    col2_x = Inches(7.0)
    col = 0

    named = [cr for cr in cluster_results if cr.cluster_id != -1]
    named.sort(key=lambda x: x.size, reverse=True)

    for i, cr in enumerate(named[:12]):
        x = col1_x if col == 0 else col2_x
        ct_str = ", ".join(f"{k}={v}" for k, v in cr.dominant_call_types.items())
        text = (
            f"● {cr.label} ({cr.size} calls)\n"
            f"   Avg sentiment: {cr.avg_sentiment_score:.1f}/5  |  {ct_str}\n"
            f"   {cr.description[:90]}"
        )
        _add_text_box(slide, text, x, y, Inches(6.2), Inches(1.0),
                      size=Pt(9.5), color=DARK)
        y += Inches(1.05)
        if y > Inches(6.8):
            col = 1
            y = Inches(1.4)


def _slide_sentiment_box(prs, transcripts, sentiment_aggregates, chart_map):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Sentiment by Call Type",
                    "Distribution and averages across all 100 transcripts")

    _add_image_safe(slide, "sentiment_by_call_type", chart_map,
                    Inches(0.3), Inches(1.3), Inches(7.5), Inches(5.5))

    # Narrative panel
    by_type = {a.call_type: a for a in sentiment_aggregates if a.cluster_label is None}
    lines = []
    for ct in ("internal", "support", "external"):
        agg = by_type.get(ct)
        if agg:
            lines.append(f"{ct.capitalize()}: avg {agg.avg_score:.1f}/5  |  "
                         f"{agg.sentence_neg_ratio*100:.0f}% negative sentences")
    interpretation = (
        "Key takeaways:\n"
        + "\n".join(f"• {l}" for l in lines)
        + "\n\nSupport calls consistently score lower than external calls — "
        "customers who reach out for help are already frustrated. "
        "Internal calls score highest, suggesting engineering teams are "
        "generally collaborative and solution-focused."
    )
    _add_text_box(slide, interpretation,
                  Inches(8.0), Inches(1.4), Inches(5.0), Inches(5.5),
                  size=Pt(11), color=DARK)


def _slide_sentiment_trend(prs, chart_map):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Sentiment Trends Over Time",
                    "Weekly average sentiment score by call type")

    _add_image_safe(slide, "sentiment_trend", chart_map,
                    Inches(0.3), Inches(1.3), Inches(8.5), Inches(5.5))

    _add_text_box(
        slide,
        "What to watch for:\n\n"
        "• Declining external sentiment over consecutive weeks = early churn signal\n\n"
        "• Support sentiment spike = potential product incident\n\n"
        "• Gap widening between internal vs external = org health concern\n\n"
        "• Sudden drop in internal sentiment = team morale or process issue",
        Inches(9.0), Inches(1.4), Inches(4.1), Inches(5.5),
        size=Pt(10.5), color=DARK,
    )


def _slide_churn_risk(prs, churn_signals: list[ChurnSignal], chart_map):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Churn Risk Analysis",
                    "Scored across external & support calls using key moments + sentiment")

    _add_image_safe(slide, "churn_risk_heatmap", chart_map,
                    Inches(0.3), Inches(1.3), Inches(5.5), Inches(3.8))

    # High risk table
    high = [c for c in churn_signals if c.risk_tier == "high"][:6]
    if high:
        y = Inches(1.4)
        _add_text_box(slide, "High-Risk Accounts:", Inches(6.1), y, Inches(6.9), Inches(0.4),
                      size=Pt(12), bold=True, color=WARN)
        y += Inches(0.45)
        for c in high:
            reasons = "; ".join(c.churn_moments[:2]) if c.churn_moments else "low sentiment"
            txt = f"• {c.title[:45]}  (score: {c.churn_score:.2f})\n  {reasons[:70]}"
            _add_text_box(slide, txt, Inches(6.1), y, Inches(6.9), Inches(0.8),
                          size=Pt(9.5), color=DARK)
            y += Inches(0.9)

    scoring_note = (
        "Scoring formula:\n"
        "churn_signal moments × 0.25 (cap 0.50)\n"
        "+ feature_gap moments × 0.10 (cap 0.25)\n"
        "+ low sentiment penalty × 0.25\n"
        "+ declining sentiment arc + 0.10\n\n"
        "High-risk calls (≥0.65) enriched by Claude claude-sonnet-4-6\n"
        "with recommended actions."
    )
    _add_text_box(slide, scoring_note, Inches(0.3), Inches(5.2), Inches(5.5), Inches(2.1),
                  size=Pt(9), color=GRAY)


def _slide_escalation(prs, escalation_flags: list[EscalationFlag], chart_map):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Escalation Detection",
                    "Multi-signal detection across all call types")

    _add_image_safe(slide, "escalation_flags", chart_map,
                    Inches(0.3), Inches(1.3), Inches(7.0), Inches(5.5))

    critical = [f for f in escalation_flags if f.severity == "critical"][:5]
    if critical:
        y = Inches(1.4)
        _add_text_box(slide, "Critical Escalations:", Inches(7.5), y, Inches(5.5), Inches(0.4),
                      size=Pt(12), bold=True, color=WARN)
        y += Inches(0.45)
        for f in critical:
            txt = (f"• {f.title[:45]}\n"
                   f"  Owner: {f.recommended_owner or 'TBD'}  |  "
                   f"{len(f.signals)} signals")
            _add_text_box(slide, txt, Inches(7.5), y, Inches(5.5), Inches(0.85),
                          size=Pt(9.5), color=DARK)
            y += Inches(0.9)


def _slide_action_items(prs, action_items, chart_map):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Action Item Ownership",
                    "Who generates the most action items and across what topics")

    _add_image_safe(slide, "action_items_treemap", chart_map,
                    Inches(0.3), Inches(1.3), Inches(6.5), Inches(5.7))

    from collections import Counter
    owner_counts = Counter(a.speaker_name for a in action_items)
    top_owners = owner_counts.most_common(8)
    lines = "\n".join(f"• {name}: {count}" for name, count in top_owners)
    _add_text_box(
        slide,
        f"Top action item owners:\n{lines}\n\n"
        f"Total action items: {len(action_items)}\n"
        f"Across {len(set(a.meeting_id for a in action_items))} calls",
        Inches(7.0), Inches(1.5), Inches(6.0), Inches(5.5),
        size=Pt(11), color=DARK,
    )


def _slide_speaker_engagement(prs, speaker_engagements, transcripts, chart_map):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Speaker Engagement Analysis",
                    "Talk-time balance and monologue risk in customer calls")

    _add_image_safe(slide, "speaker_talk_time", chart_map,
                    Inches(0.3), Inches(1.3), Inches(7.5), Inches(5.5))

    monologue_risk = [
        e for e in speaker_engagements
        if e.talk_time_pct > 70 and e.is_internal
    ]
    ct_map = {r.meeting_id: r.call_type for r in transcripts}
    external_monologue = [
        e for e in monologue_risk
        if ct_map.get(e.meeting_id) in ("external", "support")
    ]

    _add_text_box(
        slide,
        f"Calls where internal speakers dominate (>70%):\n"
        f"  Total: {len(monologue_risk)}\n"
        f"  In customer-facing calls: {len(external_monologue)}\n\n"
        "Insight: When Aegis Cloud speakers dominate customer calls, "
        "it may indicate a missed opportunity to listen. "
        "Calls with higher customer talk-time tend to score higher on sentiment — "
        "customers feel heard.\n\n"
        "Recommended: Flag calls where host talk% > 70% in external/support "
        "for coaching review.",
        Inches(8.0), Inches(1.5), Inches(5.1), Inches(5.5),
        size=Pt(10.5), color=DARK,
    )


def _slide_architecture(prs):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Pipeline Architecture & Technical Approach", "")

    pipeline = (
        "LangGraph Pipeline (8 nodes):\n\n"
        "  ingest → classify_call_type → [llm_fallback?] → embed_transcripts\n"
        "         → cluster_topics → label_clusters → analyze_sentiment\n"
        "         → extract_insights → compile_report\n"
    )
    _add_text_box(slide, pipeline, Inches(0.4), Inches(1.3), Inches(12.5), Inches(1.8),
                  size=Pt(11), color=DARK)

    decisions = (
        "Key architectural decisions:\n\n"
        "• Hybrid classification: Rule-based first (email domains + title keywords) → "
        "Claude claude-haiku-4-5 fallback for ambiguous cases. "
        "Result: fast, cheap, explainable — LLM only used where needed.\n\n"
        "• Two-pass UMAP: 10D (min_dist=0.0) for HDBSCAN clustering quality; "
        "2D (min_dist=0.1) for human-readable scatter plot visualization.\n\n"
        "• HDBSCAN over K-Means: no need to pre-specify k; handles one-off transcripts "
        "as noise (-1); finds variable-density clusters.\n\n"
        "• Local embeddings (all-MiniLM-L6-v2): 384-dim, cosine-normalized, "
        "no OpenAI dependency — all-in on Anthropic SDK.\n\n"
        "• Churn enrichment: only high-risk calls (score ≥ 0.65) sent to Claude claude-sonnet-4-6 "
        "for recommended actions — cost-efficient."
    )
    _add_text_box(slide, decisions, Inches(0.4), Inches(3.2), Inches(12.5), Inches(4.0),
                  size=Pt(10.5), color=DARK)


def _slide_recommendations(prs, churn_signals, escalation_flags, transcripts):
    slide = _blank_slide(prs)
    _add_bg(slide)
    _add_header_bar(slide, "Recommendations", "Top 3 actions from this analysis")

    high_churn = sum(1 for c in churn_signals if c.risk_tier == "high")
    critical_esc = sum(1 for e in escalation_flags if e.severity == "critical")
    support_avg = mean(
        r.sentiment_score for r in transcripts if r.call_type == "support"
    ) if any(r.call_type == "support" for r in transcripts) else 0

    recs = [
        (
            "1. Immediate account health review",
            f"There are {high_churn} high-risk accounts flagged for churn. "
            "Each has specific signals (churn_signal key moments, declining sentiment, feature gaps). "
            "Assign CS owners this week — see churn risk report for recommended actions per account.",
            WARN,
        ),
        (
            "2. Support experience audit",
            f"Support calls average {support_avg:.1f}/5 sentiment — the lowest of any call type. "
            "Cross-reference recurring topics (from clustering) with the product roadmap. "
            "Quick wins: identify top 3 support topics and add self-serve documentation.",
            ACCENT,
        ),
        (
            f"3. Escalation routing automation ({critical_esc} critical cases)",
            "The escalation detection model identified calls needing engineering, exec, or AM escalation. "
            "This logic can be productionized as a real-time webhook: flag calls during or immediately after "
            "they end, rather than in post-hoc analysis. Estimated setup: 2-3 sprints.",
            ACCENT2,
        ),
    ]

    y = Inches(1.4)
    for title, body, color in recs:
        box = slide.shapes.add_shape(1, Inches(0.4), y, Inches(12.5), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
        box.line.color.rgb = color
        box.line.width = Emu(18000)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(10)
        p2.font.color.rgb = DARK

        y += Inches(1.8)
